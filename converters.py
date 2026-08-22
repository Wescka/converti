from __future__ import annotations

import base64
import csv
import io
import html as html_lib
import json
import os
import re
import shutil
import subprocess
import xml.etree.ElementTree as ET
import tempfile
import zipfile
from pathlib import Path

from PIL import Image, ImageOps

from capabilities import AUDIO_EXTS, IMAGE_EXTS, VIDEO_EXTS
from config import MAX_IMAGE_PIXELS, MAX_PDF_PAGES, TEMP_DIR
from security_utils import Toolchain, normalize_ext, random_file, run_command

Image.MAX_IMAGE_PIXELS = MAX_IMAGE_PIXELS


def _image_save_format(ext: str) -> str:
    return {"jpg":"JPEG","png":"PNG","webp":"WEBP","bmp":"BMP","gif":"GIF","tiff":"TIFF","pdf":"PDF"}.get(ext, ext.upper())


def _rasterize_image(source: Path, tools: Toolchain, *, alpha: bool = True) -> Image.Image:
    """Open an image safely and return a detached PIL image.

    Pillow is attempted first because it avoids an external process for common
    formats. If Pillow cannot decode the source (SVG/HEIC/PSD/AVIF on some
    installations), ImageMagick is used as a decoder and the first frame/page
    is rasterized to PNG. This keeps the rest of the document pipeline
    independent of ImageMagick's PDF/Office policies.
    """
    try:
        with Image.open(source) as im:
            try:
                im.seek(0)
            except Exception:
                pass
            img = ImageOps.exif_transpose(im)
            return img.convert("RGBA" if alpha else "RGB").copy()
    except Exception as pillow_exc:
        if not tools.magick:
            raise RuntimeError(
                "Este formato de imagen necesita ImageMagick en el servidor. "
                f"Pillow no pudo abrirlo: {pillow_exc}"
            ) from pillow_exc

        tmp = random_file("png", "raster")
        # [0] fuerza la primera página/frame para PSD/GIF/ICO multipágina.
        src_arg = f"{source}[0]"
        try:
            run_command([tools.magick, src_arg, "-auto-orient", str(tmp)], timeout=90)
            with Image.open(tmp) as im:
                img = ImageOps.exif_transpose(im)
                return img.convert("RGBA" if alpha else "RGB").copy()
        finally:
            tmp.unlink(missing_ok=True)


def validate_image_source(source: Path, source_ext: str, tools: Toolchain) -> None:
    """Fail early when an uploaded image cannot be decoded safely."""
    source_ext = normalize_ext(source_ext)
    if source_ext == "svg":
        try:
            raw = source.read_text(encoding="utf-8", errors="ignore")[:2_000_000].lower()
            dangerous = ("<script", "javascript:", "file://", "http://", "https://")
            if any(token in raw for token in dangerous):
                raise ValueError("El SVG contiene referencias externas o scripts y fue rechazado por seguridad.")
        except ValueError:
            raise
        except Exception:
            pass
    img = _rasterize_image(source, tools, alpha=True)
    try:
        if img.width < 1 or img.height < 1:
            raise ValueError("La imagen no contiene dimensiones válidas.")
    finally:
        try:
            img.close()
        except Exception:
            pass


def _flatten_for_jpeg(img: Image.Image) -> Image.Image:
    if img.mode in {"RGBA", "LA"}:
        bg = Image.new("RGB", img.size, "white")
        alpha = img.getchannel("A") if "A" in img.getbands() else None
        bg.paste(img.convert("RGBA"), mask=alpha)
        return bg
    return img.convert("RGB")


def _save_raster_output(img: Image.Image, output: Path, target_ext: str, quality: int) -> None:
    target_ext = normalize_ext(target_ext)
    if target_ext == "jpg":
        img = _flatten_for_jpeg(img)
    elif target_ext in {"bmp", "pdf"} and img.mode not in {"RGB", "L"}:
        img = img.convert("RGB")
    elif target_ext == "gif":
        img = img.convert("P", palette=Image.Palette.ADAPTIVE)

    kwargs = {}
    if target_ext in {"jpg", "webp"}:
        kwargs["quality"] = quality
    img.save(output, _image_save_format(target_ext), **kwargs)


def convert_image(source: Path, target_ext: str, tools: Toolchain, options: dict) -> list[Path]:
    target_ext = normalize_ext(target_ext)
    if target_ext not in {"png", "jpg", "webp", "bmp", "gif", "tiff"}:
        raise ValueError(f"La salida .{target_ext} no es un formato de imagen directo.")

    output = random_file(target_ext, "imagen")
    quality = max(1, min(100, int(options.get("quality", 92))))

    # ImageMagick is preferred for broad input support (SVG, HEIC, PSD, AVIF),
    # but if it rejects the conversion we still try Pillow before failing.
    if tools.magick:
        try:
            cmd = [tools.magick, f"{source}[0]", "-auto-orient"]
            if target_ext in {"jpg", "webp"}:
                cmd += ["-quality", str(quality)]
            if options.get("strip_metadata"):
                cmd += ["-strip"]
            cmd.append(str(output))
            run_command(cmd, timeout=90)
            if output.exists() and output.stat().st_size > 0:
                return [output]
        except Exception:
            output.unlink(missing_ok=True)

    img = _rasterize_image(source, tools, alpha=True)
    _save_raster_output(img, output, target_ext, quality)
    return [output]


def images_to_pdf(sources: list[Path], options: dict, tools: Toolchain | None = None) -> list[Path]:
    output = random_file("pdf", "imagenes")
    tools = tools or Toolchain(None, None, None, None, None, None)
    imgs: list[Image.Image] = []
    try:
        for p in sources:
            img = _rasterize_image(p, tools, alpha=False).convert("RGB")
            imgs.append(img)
        if not imgs:
            raise ValueError("No se recibieron imágenes válidas.")
        imgs[0].save(
            output,
            "PDF",
            save_all=True,
            append_images=imgs[1:],
            resolution=float(options.get("dpi", 150)),
        )
        return [output]
    finally:
        for img in imgs:
            try:
                img.close()
            except Exception:
                pass


def _image_docx(sources: list[Path], tools: Toolchain, options: dict) -> Path:
    from docx import Document
    from docx.enum.section import WD_ORIENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Cm, Pt

    doc = Document()
    for idx, source in enumerate(sources):
        img = _rasterize_image(source, tools, alpha=False).convert("RGB")
        tmp = random_file("png", "docimg")
        try:
            img.save(tmp, "PNG")
            section = doc.sections[-1]
            wide = img.width > img.height * 1.22
            if wide:
                section.orientation = WD_ORIENT.LANDSCAPE
                section.page_width, section.page_height = section.page_height, section.page_width
            section.top_margin = Cm(1.25)
            section.bottom_margin = Cm(1.25)
            section.left_margin = Cm(1.25)
            section.right_margin = Cm(1.25)
            usable_w = section.page_width - section.left_margin - section.right_margin
            usable_h = section.page_height - section.top_margin - section.bottom_margin

            # python-docx can size by width only; cap by both dimensions using
            # the pixel aspect ratio and EMU values from the section.
            ratio = img.width / max(1, img.height)
            width_emu = int(usable_w)
            height_emu = int(width_emu / ratio)
            if height_emu > int(usable_h):
                height_emu = int(usable_h)
                width_emu = int(height_emu * ratio)

            para = doc.add_paragraph()
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            para.paragraph_format.space_before = Pt(0)
            para.paragraph_format.space_after = Pt(0)
            para.add_run().add_picture(str(tmp), width=width_emu, height=height_emu)
            if idx < len(sources) - 1:
                doc.add_page_break()
        finally:
            tmp.unlink(missing_ok=True)
            try:
                img.close()
            except Exception:
                pass

    output = random_file("docx", "imagenes")
    doc.save(output)
    return output


def _image_pptx(sources: list[Path], tools: Toolchain, options: dict) -> Path:
    from pptx import Presentation

    prs = Presentation()
    # Remove the default slide only by simply not using it; Presentation starts
    # with no slides. Use the blank layout.
    blank = prs.slide_layouts[6]
    for source in sources:
        img = _rasterize_image(source, tools, alpha=False).convert("RGB")
        tmp = random_file("png", "slideimg")
        try:
            img.save(tmp, "PNG")
            slide = prs.slides.add_slide(blank)
            sw, sh = int(prs.slide_width), int(prs.slide_height)
            ratio = img.width / max(1, img.height)
            w = sw
            h = int(w / ratio)
            if h > sh:
                h = sh
                w = int(h * ratio)
            left = int((sw - w) / 2)
            top = int((sh - h) / 2)
            slide.shapes.add_picture(str(tmp), left, top, width=w, height=h)
        finally:
            tmp.unlink(missing_ok=True)
            try:
                img.close()
            except Exception:
                pass
    output = random_file("pptx", "imagenes")
    prs.save(output)
    return output


def _image_html(sources: list[Path], tools: Toolchain, options: dict) -> Path:
    blocks = []
    for idx, source in enumerate(sources, start=1):
        img = _rasterize_image(source, tools, alpha=True)
        buf = io.BytesIO()
        try:
            img.save(buf, "PNG", optimize=True)
        finally:
            try:
                img.close()
            except Exception:
                pass
        encoded = base64.b64encode(buf.getvalue()).decode("ascii")
        blocks.append(
            f'<figure><img src="data:image/png;base64,{encoded}" alt="Imagen {idx}"></figure>'
        )
    output = random_file("html", "imagenes")
    title = html_lib.escape(str(options.get("title") or "Documento creado con Converti"))
    output.write_text(
        "<!doctype html><html lang=\"es\"><head><meta charset=\"utf-8\">"
        f"<title>{title}</title><style>body{{margin:0;background:#fff}}"
        "figure{margin:24px auto;max-width:1100px;text-align:center;page-break-after:always}"
        "figure:last-child{page-break-after:auto}img{max-width:100%;height:auto}</style></head><body>"
        + "".join(blocks) + "</body></html>",
        encoding="utf-8",
    )
    return output


def _tesseract_languages(tools: Toolchain) -> str | None:
    if not tools.tesseract:
        return None
    try:
        out = run_command([tools.tesseract, "--list-langs"], timeout=20).stdout.lower().splitlines()
        langs = {x.strip() for x in out if x.strip() and not x.lower().startswith("list of available")}
        wanted = [x for x in ("spa", "eng", "por", "fra") if x in langs]
        return "+".join(wanted) if wanted else None
    except Exception:
        return None


def _image_ocr_txt(sources: list[Path], tools: Toolchain, options: dict) -> Path:
    if not tools.tesseract:
        raise RuntimeError("Tesseract OCR no está disponible para convertir la imagen a texto.")
    lang = _tesseract_languages(tools)
    pages = []
    for source in sources:
        img = _rasterize_image(source, tools, alpha=False).convert("RGB")
        tmp = random_file("png", "ocr")
        try:
            img.save(tmp, "PNG")
            cmd = [tools.tesseract, str(tmp), "stdout", "--psm", "6"]
            if lang:
                cmd += ["-l", lang]
            pages.append(run_command(cmd, timeout=120).stdout.strip())
        finally:
            tmp.unlink(missing_ok=True)
            try:
                img.close()
            except Exception:
                pass
    output = random_file("txt", "ocr")
    output.write_text("\n\n--- Página siguiente ---\n\n".join(pages), encoding="utf-8")
    return output


def images_to_document(sources: list[Path], target_ext: str, tools: Toolchain, options: dict) -> list[Path]:
    """Combine one or more images into a real document container."""
    target_ext = normalize_ext(target_ext)
    if not sources:
        raise ValueError("No se recibieron imágenes válidas.")

    if target_ext == "pdf":
        return images_to_pdf(sources, options, tools)
    if target_ext == "docx":
        return [_image_docx(sources, tools, options)]
    if target_ext == "pptx":
        return [_image_pptx(sources, tools, options)]
    if target_ext == "html":
        return [_image_html(sources, tools, options)]
    if target_ext == "txt":
        return [_image_ocr_txt(sources, tools, options)]
    if target_ext in {"odt", "rtf"}:
        if not tools.soffice:
            raise RuntimeError("LibreOffice no está disponible para generar este documento.")
        intermediate = _image_docx(sources, tools, options)
        try:
            return convert_with_soffice(intermediate, target_ext, tools)
        finally:
            intermediate.unlink(missing_ok=True)
    raise ValueError(f"Imagen a .{target_ext} no está implementado.")

def _probe_duration_seconds(source: Path, tools: Toolchain) -> float | None:
    if not tools.ffprobe:
        return None
    try:
        result = subprocess.run(
            [tools.ffprobe, "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", str(source)],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=20, creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        value = float((result.stdout or "").strip())
        return value if value > 0 else None
    except Exception:
        return None


def convert_media(source: Path, target_ext: str, tools: Toolchain, options: dict, progress_callback=None) -> list[Path]:
    if not tools.ffmpeg:
        raise RuntimeError("FFmpeg no está disponible.")
    target_ext = normalize_ext(target_ext)
    output = random_file(target_ext, "media")
    bitrate = str(options.get("audio_bitrate", "192k"))
    resolution = str(options.get("resolution", "original"))
    quality = max(18, min(35, int(options.get("video_crf", 23))))
    duration = _probe_duration_seconds(source, tools)
    cmd = [tools.ffmpeg, "-y", "-hide_banner", "-loglevel", "error", "-progress", "pipe:1", "-nostats", "-i", str(source)]
    if target_ext in AUDIO_EXTS:
        cmd += ["-vn"]
        codec_map = {"mp3":"libmp3lame","flac":"flac","ogg":"libvorbis","opus":"libopus","aac":"aac","m4a":"aac","wav":"pcm_s16le"}
        if target_ext in codec_map:
            cmd += ["-c:a", codec_map[target_ext]]
        if target_ext not in {"flac","wav"}:
            cmd += ["-b:a", bitrate]
    else:
        if resolution != "original" and "x" in resolution:
            cmd += ["-vf", f"scale={resolution.replace('x',':')}:force_original_aspect_ratio=decrease"]
        if target_ext in {"mp4","mkv","mov"}:
            cmd += ["-c:v","libx264","-crf",str(quality),"-c:a","aac"]
        elif target_ext == "webm":
            cmd += ["-c:v","libvpx-vp9","-crf",str(quality),"-b:v","0","-c:a","libopus"]
    if options.get("strip_metadata"):
        cmd += ["-map_metadata", "-1"]
    cmd.append(str(output))
    proc = subprocess.Popen(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, bufsize=1, universal_newlines=True,
        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
    )
    try:
        if proc.stdout is not None:
            for raw in proc.stdout:
                line = raw.strip()
                if progress_callback and duration and line.startswith("out_time_ms="):
                    try:
                        processed = int(line.split("=", 1)[1]) / 1_000_000.0
                        pct = max(0.0, min(99.0, processed / duration * 100.0))
                        progress_callback(pct)
                    except Exception:
                        pass
                elif progress_callback and line == "progress=end":
                    progress_callback(100.0)
        stderr = proc.stderr.read() if proc.stderr is not None else ""
        code = proc.wait()
        if code != 0:
            raise RuntimeError((stderr or "FFmpeg no pudo completar la conversión.")[-4000:])
    finally:
        if proc.poll() is None:
            proc.kill()
    if progress_callback:
        progress_callback(100.0)
    return [output]


def convert_with_soffice(source: Path, target_ext: str, tools: Toolchain) -> list[Path]:
    """LibreOffice conversion with a Termux/proot-safe path bridge.

    On the Redmi, `soffice` is a wrapper that enters Debian with --shared-tmp.
    Host $PREFIX/tmp is visible as /tmp inside Debian, so paths must be translated.
    """
    if not tools.soffice:
        raise RuntimeError("LibreOffice no está disponible.")

    prefix = os.environ.get("PREFIX", "")
    is_termux = bool(prefix and "com.termux" in prefix)

    if is_termux:
        shared_tmp = Path(prefix) / "tmp"
        shared_tmp.mkdir(parents=True, exist_ok=True)
        job = Path(tempfile.mkdtemp(prefix="converti_lo_", dir=shared_tmp))
        out_host = job / "out"
        profile_host = job / "profile"
        out_host.mkdir()
        profile_host.mkdir()
        # Copy into shared tmp so Debian can read it through /tmp.
        src_host = job / ("source." + normalize_ext(source.suffix.lstrip(".")))
        shutil.copy2(source, src_host)
        guest_root = f"/tmp/{job.name}"
        guest_src = f"{guest_root}/{src_host.name}"
        guest_out = f"{guest_root}/out"
        guest_profile_uri = f"file://{guest_root}/profile"
        try:
            run_command([
                tools.soffice,
                "--headless", "--nologo", "--nodefault", "--nolockcheck",
                f"-env:UserInstallation={guest_profile_uri}",
                "--convert-to", target_ext,
                "--outdir", guest_out,
                guest_src,
            ], timeout=120)
            candidates = [x for x in out_host.iterdir() if x.is_file()]
            if not candidates:
                raise RuntimeError("LibreOffice no generó el archivo de salida.")
            produced = max(candidates, key=lambda x: x.stat().st_mtime)
            output = random_file(target_ext, "office")
            shutil.move(str(produced), output)
            return [output]
        finally:
            shutil.rmtree(job, ignore_errors=True)

    work = Path(tempfile.mkdtemp(prefix="converti_lo_", dir=TEMP_DIR))
    profile = Path(tempfile.mkdtemp(prefix="converti_profile_", dir=TEMP_DIR))
    try:
        profile_uri = profile.resolve().as_uri()
        run_command([
            tools.soffice,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            f"-env:UserInstallation={profile_uri}",
            "--convert-to", target_ext,
            "--outdir", str(work),
            str(source),
        ], timeout=120)
        candidates = [x for x in work.iterdir() if x.is_file()]
        if not candidates:
            raise RuntimeError("LibreOffice no generó el archivo de salida.")
        produced = max(candidates, key=lambda x: x.stat().st_mtime)
        output = random_file(target_ext, "office")
        shutil.move(str(produced), output)
        return [output]
    finally:
        shutil.rmtree(work, ignore_errors=True)
        shutil.rmtree(profile, ignore_errors=True)


def convert_with_pandoc(source: Path, target_ext: str, tools: Toolchain) -> list[Path]:
    if not tools.pandoc:
        raise RuntimeError("Pandoc no está disponible.")
    output = random_file(target_ext, "pandoc")
    target_map = {"md":"gfm","txt":"plain","html":"html5","tex":"latex"}
    cmd = [tools.pandoc, str(source), "-t", target_map.get(target_ext, target_ext), "-o", str(output)]
    run_command(cmd)
    return [output]


def _has_pymupdf() -> bool:
    try:
        import fitz  # noqa: F401
        return True
    except Exception:
        try:
            import pymupdf  # noqa: F401
            return True
        except Exception:
            return False


def _mutool() -> str | None:
    return shutil.which("mutool")


def _pdf_reader(source: Path):
    from pypdf import PdfReader
    return PdfReader(str(source))


def _pdf_page_count(source: Path) -> int:
    try:
        return len(_pdf_reader(source).pages)
    except Exception as exc:
        raise RuntimeError(f"No se pudo leer el PDF: {exc}") from exc


def _selected_page_ids(total: int, options: dict) -> list[int]:
    return list(_page_indexes(total, options))


def _extract_page_texts_pypdf(source: Path, page_ids: list[int]) -> list[str]:
    reader = _pdf_reader(source)
    texts = []
    for i in page_ids:
        try:
            texts.append((reader.pages[i].extract_text() or "").strip())
        except Exception:
            texts.append("")
    return texts


def _render_page_with_mutool(source: Path, page_number_1based: int, dpi: int = 150) -> Path:
    mutool = _mutool()
    if not mutool:
        raise RuntimeError("MuPDF (mutool) no está disponible.")
    out = random_file("png", f"pagina_{page_number_1based}")
    cmd = [
        mutool, "draw",
        "-q",
        "-r", str(dpi),
        "-F", "png",
        "-o", str(out),
        str(source),
        str(page_number_1based),
    ]
    run_command(cmd)
    if not out.exists() or out.stat().st_size == 0:
        raise RuntimeError("MuPDF no generó la imagen de la página.")
    return out


def pdf_to_textual(source: Path, target_ext: str, tools: Toolchain, options: dict) -> list[Path]:
    """
    Usa PyMuPDF cuando está disponible (PC/Render).
    En Android/Termux cae automáticamente a pypdf + mutool.
    """
    from docx import Document

    # Ruta original: conserva mejor imágenes/bloques cuando PyMuPDF existe.
    try:
        import fitz
    except Exception:
        fitz = None

    if fitz is not None:
        doc = fitz.open(source)
        if len(doc) > MAX_PDF_PAGES:
            raise ValueError(f"El PDF supera el máximo de {MAX_PDF_PAGES} páginas permitido.")
        pages_text = []
        page_ids = list(_page_indexes(len(doc), options))
        for i in page_ids:
            page = doc[i]
            text = page.get_text("text").strip()
            if not text and options.get("ocr") and tools.tesseract:
                pix = page.get_pixmap(matrix=fitz.Matrix(2,2), alpha=False)
                temp_png = random_file("png", f"ocr_{i}")
                pix.save(temp_png)
                outbase = temp_png.with_suffix("")
                run_command([tools.tesseract, str(temp_png), str(outbase), "-l", str(options.get("ocr_lang", "spa+eng"))])
                txt_path = outbase.with_suffix(".txt")
                text = txt_path.read_text(encoding="utf-8", errors="ignore") if txt_path.exists() else ""
                temp_png.unlink(missing_ok=True)
                txt_path.unlink(missing_ok=True)
            pages_text.append(text)

        if target_ext == "txt":
            out = random_file("txt", "pdf")
            out.write_text("\n\n".join(pages_text), encoding="utf-8")
            return [out]
        if target_ext == "md":
            out = random_file("md", "pdf")
            out.write_text("\n\n---\n\n".join(pages_text), encoding="utf-8")
            return [out]
        if target_ext == "html":
            out = random_file("html", "pdf")
            body = "\n".join(f"<section><pre>{_escape(t)}</pre></section>" for t in pages_text)
            out.write_text(f"<!doctype html><meta charset='utf-8'><body>{body}</body>", encoding="utf-8")
            return [out]
        if target_ext == "docx":
            from io import BytesIO
            from docx.shared import Inches
            out = random_file("docx", "pdf")
            word = Document()
            for seq, page_index in enumerate(page_ids):
                page = doc[page_index]
                if seq > 0:
                    word.add_page_break()
                data = page.get_text("dict")
                wrote = False
                for block in data.get("blocks", []):
                    if block.get("type") == 0:
                        lines = []
                        for line in block.get("lines", []):
                            text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
                            if text:
                                lines.append(text)
                        if lines:
                            word.add_paragraph("\n".join(lines))
                            wrote = True
                    elif block.get("type") == 1 and block.get("image"):
                        try:
                            stream = BytesIO(block["image"])
                            word.add_picture(stream, width=Inches(5.8))
                            wrote = True
                        except Exception:
                            pass
                if not wrote:
                    fallback = pages_text[seq]
                    for para in [x.strip() for x in fallback.split("\n") if x.strip()]:
                        word.add_paragraph(para)
            word.save(out)
            return [out]
        raise ValueError("Destino textual de PDF no soportado.")

    # Fallback Termux/Android: pypdf para texto y mutool para OCR/render.
    total = _pdf_page_count(source)
    if total > MAX_PDF_PAGES:
        raise ValueError(f"El PDF supera el máximo de {MAX_PDF_PAGES} páginas permitido.")
    page_ids = _selected_page_ids(total, options)
    pages_text = _extract_page_texts_pypdf(source, page_ids)

    if options.get("ocr") and tools.tesseract:
        for seq, page_index in enumerate(page_ids):
            if pages_text[seq]:
                continue
            temp_png = None
            txt_path = None
            try:
                temp_png = _render_page_with_mutool(source, page_index + 1, 200)
                outbase = temp_png.with_suffix("")
                run_command([
                    tools.tesseract, str(temp_png), str(outbase),
                    "-l", str(options.get("ocr_lang", "spa+eng"))
                ])
                txt_path = outbase.with_suffix(".txt")
                if txt_path.exists():
                    pages_text[seq] = txt_path.read_text(encoding="utf-8", errors="ignore").strip()
            finally:
                if temp_png:
                    temp_png.unlink(missing_ok=True)
                if txt_path:
                    txt_path.unlink(missing_ok=True)

    if target_ext == "txt":
        out = random_file("txt", "pdf")
        out.write_text("\n\n".join(pages_text), encoding="utf-8")
        return [out]

    if target_ext == "md":
        out = random_file("md", "pdf")
        out.write_text("\n\n---\n\n".join(pages_text), encoding="utf-8")
        return [out]

    if target_ext == "html":
        out = random_file("html", "pdf")
        body = "\n".join(f"<section><pre>{_escape(t)}</pre></section>" for t in pages_text)
        out.write_text(f"<!doctype html><meta charset='utf-8'><body>{body}</body>", encoding="utf-8")
        return [out]

    if target_ext == "docx":
        out = random_file("docx", "pdf")
        word = Document()
        for seq, text in enumerate(pages_text):
            if seq > 0:
                word.add_page_break()
            paragraphs = [x.strip() for x in text.splitlines() if x.strip()]
            if paragraphs:
                for paragraph in paragraphs:
                    word.add_paragraph(paragraph)
            else:
                word.add_paragraph("[Página sin texto extraíble]")
        word.save(out)
        return [out]

    raise ValueError("Destino textual de PDF no soportado.")


def _escape(value: str) -> str:
    import html
    return html.escape(value)


def _page_indexes(total: int, options: dict) -> range:
    start = max(1, int(options.get("page_start", 1) or 1))
    end_raw = int(options.get("page_end", 0) or 0)
    end = total if end_raw <= 0 else min(total, end_raw)
    if start > end:
        raise ValueError("El rango de páginas no es válido.")
    return range(start - 1, end)


def pdf_to_images(source: Path, target_ext: str, options: dict) -> list[Path]:
    """
    PyMuPDF en sistemas donde esté instalado.
    MuPDF CLI (mutool) como fallback nativo de Termux.
    """
    try:
        import fitz
    except Exception:
        fitz = None

    dpi = max(72, min(600, int(options.get("dpi", 150))))
    quality = max(1, min(100, int(options.get("quality", 92))))

    if fitz is not None:
        doc = fitz.open(source)
        if len(doc) > MAX_PDF_PAGES:
            raise ValueError(f"El PDF supera el máximo de {MAX_PDF_PAGES} páginas permitido.")
        scale = dpi / 72
        outputs = []
        for idx in _page_indexes(len(doc), options):
            page = doc[idx]
            pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            if target_ext == "png":
                out = random_file("png", f"pagina_{idx+1}")
                pix.save(out)
            else:
                tmp = random_file("png", f"pagina_{idx+1}")
                pix.save(tmp)
                out = random_file("jpg", f"pagina_{idx+1}")
                with Image.open(tmp) as im:
                    im.convert("RGB").save(out, "JPEG", quality=quality)
                tmp.unlink(missing_ok=True)
            outputs.append(out)
        return outputs

    total = _pdf_page_count(source)
    if total > MAX_PDF_PAGES:
        raise ValueError(f"El PDF supera el máximo de {MAX_PDF_PAGES} páginas permitido.")
    if not _mutool():
        raise RuntimeError("PDF a imagen requiere PyMuPDF o MuPDF (mutool).")

    outputs = []
    for idx in _page_indexes(total, options):
        tmp = _render_page_with_mutool(source, idx + 1, dpi)
        if target_ext == "png":
            out = random_file("png", f"pagina_{idx+1}")
            shutil.move(str(tmp), str(out))
        else:
            out = random_file("jpg", f"pagina_{idx+1}")
            with Image.open(tmp) as im:
                im.convert("RGB").save(out, "JPEG", quality=quality)
            tmp.unlink(missing_ok=True)
        outputs.append(out)
    return outputs


def convert_structured(source: Path, source_ext: str, target_ext: str) -> list[Path]:
    import pandas as pd
    source_ext = normalize_ext(source_ext)
    target_ext = normalize_ext(target_ext)
    if source_ext == "csv":
        try:
            df = pd.read_csv(source)
        except UnicodeDecodeError:
            df = pd.read_csv(source, encoding="latin-1")
    elif source_ext in {"xlsx", "xls", "ods"}:
        df = pd.read_excel(source)
    elif source_ext == "json":
        data = json.loads(source.read_text(encoding="utf-8", errors="ignore"))
        df = pd.json_normalize(data if isinstance(data, list) else [data])
    elif source_ext == "xml":
        root = ET.parse(source).getroot()
        rows = []
        for child in list(root):
            row = {sub.tag: (sub.text or "") for sub in list(child)}
            if not row:
                row = {child.tag: child.text or ""}
            rows.append(row)
        df = pd.DataFrame(rows)
    else:
        raise ValueError("Formato de datos no soportado.")
    out = random_file(target_ext, "datos")
    if target_ext == "csv":
        df.to_csv(out, index=False)
    elif target_ext == "xlsx":
        df.to_excel(out, index=False)
    elif target_ext == "json":
        df.to_json(out, orient="records", indent=2, force_ascii=False)
    elif target_ext == "txt":
        out.write_text(df.to_string(index=False), encoding="utf-8")
    elif target_ext == "xml":
        root = ET.Element("rows")
        for record in df.fillna("").to_dict(orient="records"):
            row_el = ET.SubElement(root, "row")
            for key, value in record.items():
                el = ET.SubElement(row_el, str(key))
                el.text = str(value)
        ET.ElementTree(root).write(out, encoding="utf-8", xml_declaration=True)
    else:
        raise ValueError("Destino de datos no soportado.")
    return [out]



def convert_docx_native(source: Path, target_ext: str, options: dict | None = None) -> list[Path]:
    """Conversor DOCX sin LibreOffice/Pandoc, pensado para Termux.

    Salidas fiables: TXT, MD, HTML y PDF. Conserva texto, títulos y tablas;
    PDF mantiene el contenido textual y una maquetación limpia.
    """
    from docx import Document
    from xml.sax.saxutils import escape as xml_escape

    target_ext = normalize_ext(target_ext)
    if target_ext not in {"txt", "md", "html", "pdf"}:
        raise ValueError(f"DOCX a .{target_ext} no está implementado por el motor nativo.")

    doc = Document(source)
    blocks: list[tuple[str, object]] = []

    # Recorrido en orden aproximado de párrafos y tablas.
    for child in doc.element.body.iterchildren():
        tag = child.tag.rsplit('}', 1)[-1]
        if tag == 'p':
            from docx.text.paragraph import Paragraph
            p = Paragraph(child, doc)
            text = (p.text or '').strip()
            if text:
                style = (p.style.name if p.style else '').lower()
                level = 0
                if 'title' in style:
                    level = 1
                elif 'heading' in style:
                    m = re.search(r'(\d+)', style)
                    level = int(m.group(1)) if m else 2
                blocks.append(('heading' if level else 'p', (level, text)))
        elif tag == 'tbl':
            from docx.table import Table
            table = Table(child, doc)
            rows = [[(c.text or '').strip() for c in row.cells] for row in table.rows]
            if rows:
                blocks.append(('table', rows))

    if target_ext == 'txt':
        out = random_file('txt', 'docx')
        lines = []
        for kind, payload in blocks:
            if kind in {'p', 'heading'}:
                _, text = payload
                lines.append(text)
            else:
                for row in payload:
                    lines.append('\t'.join(row))
            lines.append('')
        out.write_text('\n'.join(lines).strip() + '\n', encoding='utf-8')
        return [out]

    if target_ext == 'md':
        out = random_file('md', 'docx')
        parts = []
        for kind, payload in blocks:
            if kind == 'heading':
                level, text = payload
                level = max(1, min(6, level or 2))
                parts.append('#' * level + ' ' + text)
            elif kind == 'p':
                _, text = payload
                parts.append(text)
            else:
                rows = payload
                if rows:
                    width = max(len(r) for r in rows)
                    first = rows[0] + [''] * (width - len(rows[0]))
                    parts.append('| ' + ' | '.join(first) + ' |')
                    parts.append('| ' + ' | '.join(['---'] * width) + ' |')
                    for row in rows[1:]:
                        row = row + [''] * (width - len(row))
                        parts.append('| ' + ' | '.join(row) + ' |')
            parts.append('')
        out.write_text('\n'.join(parts).strip() + '\n', encoding='utf-8')
        return [out]

    if target_ext == 'html':
        out = random_file('html', 'docx')
        parts = ["<!doctype html><html><head><meta charset='utf-8'>",
                 "<meta name='viewport' content='width=device-width,initial-scale=1'>",
                 "<style>body{font:16px/1.55 Arial,sans-serif;max-width:900px;margin:40px auto;padding:0 20px;color:#17213b}table{border-collapse:collapse;width:100%;margin:18px 0}td,th{border:1px solid #ccd5e5;padding:8px;text-align:left}</style></head><body>"]
        for kind, payload in blocks:
            if kind == 'heading':
                level, text = payload
                level = max(1, min(6, level or 2))
                parts.append(f'<h{level}>{xml_escape(text)}</h{level}>')
            elif kind == 'p':
                _, text = payload
                parts.append(f'<p>{xml_escape(text)}</p>')
            else:
                rows = payload
                parts.append('<table>')
                for ridx, row in enumerate(rows):
                    cell = 'th' if ridx == 0 else 'td'
                    parts.append('<tr>' + ''.join(f'<{cell}>{xml_escape(v)}</{cell}>' for v in row) + '</tr>')
                parts.append('</table>')
        parts.append('</body></html>')
        out.write_text(''.join(parts), encoding='utf-8')
        return [out]

    # DOCX -> PDF con ReportLab, sin LibreOffice.
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib import colors
    from reportlab.lib.units import mm

    out = random_file('pdf', 'docx')
    pdf = SimpleDocTemplate(str(out), pagesize=A4, rightMargin=18*mm, leftMargin=18*mm, topMargin=18*mm, bottomMargin=18*mm)
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle('ConvertiBody', parent=styles['BodyText'], fontName='Helvetica', fontSize=10.5, leading=15, spaceAfter=7)
    story = []
    for kind, payload in blocks:
        if kind == 'heading':
            level, text = payload
            style_name = 'Title' if level == 1 else ('Heading1' if level <= 2 else 'Heading2')
            story.append(Paragraph(xml_escape(text), styles[style_name]))
            story.append(Spacer(1, 4))
        elif kind == 'p':
            _, text = payload
            story.append(Paragraph(xml_escape(text).replace('\n', '<br/>'), body_style))
        else:
            rows = payload
            if rows:
                safe_rows = [[Paragraph(xml_escape(v), body_style) for v in row] for row in rows]
                table = Table(safe_rows, repeatRows=1, hAlign='LEFT')
                table.setStyle(TableStyle([
                    ('GRID', (0,0), (-1,-1), 0.4, colors.HexColor('#cbd5e1')),
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f1f5f9')),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ('LEFTPADDING', (0,0), (-1,-1), 5),
                    ('RIGHTPADDING', (0,0), (-1,-1), 5),
                    ('TOPPADDING', (0,0), (-1,-1), 5),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 5),
                ]))
                story.append(table)
                story.append(Spacer(1, 8))
    if not story:
        story.append(Paragraph('Documento sin texto extraíble.', body_style))
    pdf.build(story)
    return [out]

def package_outputs(outputs: list[Path], download_base: str) -> tuple[Path, str]:
    if len(outputs) == 1:
        return outputs[0], f"{download_base}.{outputs[0].suffix.lstrip('.')}"
    zip_path = random_file("zip", "resultado")
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
        for i, p in enumerate(outputs, start=1):
            z.write(p, arcname=f"{download_base}_{i}{p.suffix}")
    for p in outputs:
        p.unlink(missing_ok=True)
    return zip_path, f"{download_base}.zip"
